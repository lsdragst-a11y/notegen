"""
课程结题展示 PPT 生成脚本（14 页，约 10 分钟）。
运行：python build_pptx.py
输出：notegen_presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- 视觉主题 ---
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x68)   # 深蓝
COLOR_ACCENT  = RGBColor(0xE8, 0x7A, 0x2E)   # 强调橙
COLOR_TEXT    = RGBColor(0x22, 0x22, 0x22)   # 近黑
COLOR_MUTED   = RGBColor(0x66, 0x66, 0x66)   # 灰
COLOR_BG_SOFT = RGBColor(0xF4, 0xF6, 0xFA)   # 浅蓝灰
COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
FONT_CN = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# --- 辅助函数 ---
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=COLOR_TEXT,
             align=PP_ALIGN.LEFT, font=FONT_CN, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=20, line_space=1.35,
                color=COLOR_TEXT, bullet_color=COLOR_ACCENT):
    """items: list of str or (str, str) -> (要点, 大白话解释)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_space
        # 圆点
        r0 = p.add_run()
        r0.text = "● "
        r0.font.name = FONT_CN
        r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color
        # 正文
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run()
            r1.text = head
            r1.font.name = FONT_CN
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = "  " + tail
            r2.font.name = FONT_CN
            r2.font.size = Pt(size - 2)
            r2.font.color.rgb = COLOR_MUTED
        else:
            r1 = p.add_run()
            r1.text = item
            r1.font.name = FONT_CN
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
    return tb

def add_header(slide, page_no, total, title):
    # 顶部条
    add_rect(slide, 0, 0, SW, Inches(0.85), COLOR_PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11), Inches(0.5),
             title, size=24, bold=True, color=COLOR_WHITE)
    add_text(slide, Inches(11.7), Inches(0.22), Inches(1.4), Inches(0.4),
             f"{page_no} / {total}", size=14, color=COLOR_WHITE,
             align=PP_ALIGN.RIGHT)

def add_footer(slide):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(12), Inches(0.35),
             "网课视频笔记生成系统 · 课程结题展示",
             size=10, color=COLOR_MUTED)

def add_placeholder_box(slide, x, y, w, h, text):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = COLOR_BG_SOFT
    shp.line.color.rgb = COLOR_MUTED
    shp.line.width = Pt(1.0)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_CN
    r.font.size = Pt(16)
    r.font.color.rgb = COLOR_MUTED
    return shp


def add_fitted_picture(slide, x, y, box_w, box_h, path, *, bg=True):
    """把图片按比例 fit 到 (x,y,box_w,box_h) 框里，居中。
    bg=True 会先在框背景画一个浅色底（让画面感更整）。"""
    from PIL import Image
    if bg:
        add_rect(slide, x, y, box_w, box_h, COLOR_BG_SOFT)
    iw, ih = Image.open(path).size
    img_ratio = iw / ih
    box_ratio = box_w / box_h
    if img_ratio > box_ratio:
        actual_w = box_w
        actual_h = int(box_w / img_ratio)
    else:
        actual_h = box_h
        actual_w = int(box_h * img_ratio)
    offset_x = x + (box_w - actual_w) // 2
    offset_y = y + (box_h - actual_h) // 2
    slide.shapes.add_picture(path, offset_x, offset_y,
                              width=actual_w, height=actual_h)

TOTAL = 14


# ============================================================
# Page 1 - 封面
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, COLOR_PRIMARY)
# 左侧装饰条
add_rect(s, Inches(0.6), Inches(2.2), Inches(0.12), Inches(3.1), COLOR_ACCENT)
add_text(s, Inches(0.95), Inches(2.15), Inches(12), Inches(0.6),
         "课程结题展示", size=18, color=RGBColor(0xCF, 0xD8, 0xE6))
add_text(s, Inches(0.95), Inches(2.7), Inches(12), Inches(1.3),
         "基于深度学习的", size=40, bold=True, color=COLOR_WHITE)
add_text(s, Inches(0.95), Inches(3.45), Inches(12), Inches(1.3),
         "网课视频摘要与笔记生成系统", size=40, bold=True, color=COLOR_WHITE)
add_text(s, Inches(0.95), Inches(4.45), Inches(12), Inches(0.6),
         "从视频到结构化笔记的端到端 Pipeline",
         size=20, color=RGBColor(0xCF, 0xD8, 0xE6))
# 信息
add_text(s, Inches(0.95), Inches(5.7), Inches(12), Inches(0.5),
         "汇报人：[姓名]    学号：[学号]",
         size=16, color=COLOR_WHITE)
add_text(s, Inches(0.95), Inches(6.15), Inches(12), Inches(0.5),
         "指导老师：[老师姓名]    日期：2026-05",
         size=16, color=RGBColor(0xCF, 0xD8, 0xE6))


# ============================================================
# Page 2 - 背景与动机
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 2, TOTAL, "一、背景与动机")
add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "为什么要做这个系统？", size=22, bold=True, color=COLOR_PRIMARY)

add_bullets(s, Inches(0.6), Inches(2.0), Inches(12), Inches(4.5), [
    ("痛点 1：网课视频太长。",
     "一节课动辄 1~2 小时，看完不可能再翻第二遍。"),
    ("痛点 2：现有自动笔记产出是一大坨文字。",
     "没有章节、没有目录、没有重点，等于把字幕复制了一遍。"),
    ("痛点 3：通用模板套不同视频都一个样。",
     "教学课和 Vlog 排版完全一样，看起来很突兀。"),
    ("目标：做一个能自动看完视频，并按学习场景排版的笔记工具。",
     "输入一个 B 站链接，输出可阅读的结构化 Markdown 笔记。"),
], size=20)
add_footer(s)


# ============================================================
# Page 3 - 任务定义 + 系统总览
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 3, TOTAL, "二、任务定义 与 系统总览")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "输入一个视频链接，输出一份结构化笔记 + 网页可视化",
         size=18, bold=True, color=COLOR_PRIMARY)

# Pipeline 流程图（7 个圆角矩形 + 箭头）
steps = ["下载视频", "语音识别\nASR", "视频分类", "章节切分\nLLM",
         "写章标题", "套模板", "md / Web\n双输出"]
box_w, box_h = Inches(1.55), Inches(1.05)
gap = Inches(0.25)
total_w = box_w * 7 + gap * 6
start_x = (SW - total_w) / 2
y = Inches(2.4)

for i, label in enumerate(steps):
    x = start_x + (box_w + gap) * i
    fill = COLOR_PRIMARY if i in (1, 3) else COLOR_BG_SOFT  # 关键模块用主色
    text_color = COLOR_WHITE if i in (1, 3) else COLOR_TEXT
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = COLOR_PRIMARY
    shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT_CN; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = text_color
    if i < len(steps) - 1:
        arrow_x = x + box_w
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 arrow_x, y + Inches(0.42),
                                 gap, Inches(0.2))
        arr.fill.solid(); arr.fill.fore_color.rgb = COLOR_ACCENT
        arr.line.fill.background()

# 说明文字
add_text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.4),
         "输入：B 站链接 或 本地视频文件",
         size=18, color=COLOR_TEXT)
add_text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
         "输出：① 结构化 Markdown 笔记    ② Next.js 网页可视化（带播放器和章节跳转）",
         size=18, color=COLOR_TEXT)

add_rect(s, Inches(0.6), Inches(5.3), Inches(12.13), Inches(1.4), COLOR_BG_SOFT)
add_text(s, Inches(0.85), Inches(5.5), Inches(11.7), Inches(0.5),
         "关键设计：全程本地运行，不依赖在线 API。",
         size=18, bold=True, color=COLOR_PRIMARY)
add_text(s, Inches(0.85), Inches(6.05), Inches(11.7), Inches(0.5),
         "所用模型：faster-whisper（语音）+ Qwen2.5-7B-AWQ（大模型）+ Qwen2.5-VL（看图）。",
         size=15, color=COLOR_MUTED)
add_footer(s)


# ============================================================
# Page 4 - 代码结构
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 4, TOTAL, "三、代码结构")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "后端 src/ 共 20+ 模块，前端 web/ 用 Next.js 16",
         size=18, bold=True, color=COLOR_PRIMARY)

# 左：后端模块表
add_text(s, Inches(0.6), Inches(1.9), Inches(6), Inches(0.4),
         "后端 src/  （Python）", size=18, bold=True, color=COLOR_TEXT)
backend = [
    ("download.py", "yt-dlp 下载 B 站视频"),
    ("asr.py", "faster-whisper 语音识别"),
    ("classify_category.py", "4 类启发式视频分类"),
    ("segment_llm.py", "LLM 章节切分（核心）"),
    ("summarize.py", "章标题 + 摘要 + 章末小结"),
    ("caption_vl.py", "视觉大模型描述画面"),
    ("keyframe.py", "关键帧抽取"),
    ("pipeline.py", "总调度，串联以上模块"),
]
table_y = Inches(2.35)
for i, (fn, desc) in enumerate(backend):
    y = table_y + Inches(0.42) * i
    add_text(s, Inches(0.7), y, Inches(2.5), Inches(0.4),
             fn, size=13, bold=True, color=COLOR_ACCENT, font="Consolas")
    add_text(s, Inches(3.25), y, Inches(3.5), Inches(0.4),
             desc, size=13, color=COLOR_TEXT)

# 右：前端 + 模型
add_text(s, Inches(7.2), Inches(1.9), Inches(6), Inches(0.4),
         "前端 web/  （TypeScript）", size=18, bold=True, color=COLOR_TEXT)
frontend = [
    ("Next.js 16 + React", "App Router + SSG"),
    ("Plyr 播放器", "视频时间戳跳转"),
    ("双语 toggle", "中文 ↔ 英文一键切换"),
    ("Apple 风格 UI", "流动粒子背景"),
]
for i, (fn, desc) in enumerate(frontend):
    y = Inches(2.35) + Inches(0.42) * i
    add_text(s, Inches(7.3), y, Inches(2.5), Inches(0.4),
             fn, size=13, bold=True, color=COLOR_ACCENT)
    add_text(s, Inches(9.85), y, Inches(3.0), Inches(0.4),
             desc, size=13, color=COLOR_TEXT)

# 底部 stats
add_rect(s, Inches(7.2), Inches(4.5), Inches(5.5), Inches(2.0), COLOR_BG_SOFT)
add_text(s, Inches(7.4), Inches(4.65), Inches(5.2), Inches(0.4),
         "运行环境", size=15, bold=True, color=COLOR_PRIMARY)
add_text(s, Inches(7.4), Inches(5.1), Inches(5.2), Inches(0.4),
         "● 显存约 5 GB（4-bit 量化）", size=13, color=COLOR_TEXT)
add_text(s, Inches(7.4), Inches(5.5), Inches(5.2), Inches(0.4),
         "● 完整 pipeline 单视频 5-20 分钟", size=13, color=COLOR_TEXT)
add_text(s, Inches(7.4), Inches(5.9), Inches(5.2), Inches(0.4),
         "● 全本地，无在线 API 依赖", size=13, color=COLOR_TEXT)
add_footer(s)


# ============================================================
# Page 5 - 方法 ① LLM 章节切分
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 5, TOTAL, "四、核心方法 ① · LLM 章节切分")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "用大模型替代传统算法切章节，配三层程序兜底",
         size=18, bold=True, color=COLOR_PRIMARY)

add_bullets(s, Inches(0.6), Inches(2.0), Inches(12), Inches(3.0), [
    ("使用本地大模型 Qwen2.5-7B-AWQ。",
     "4-bit 量化版本，5 GB 显存就能跑。"),
    ("B1 两步法：先列大纲，再起标题。",
     "拆成两次调用，避免一次性出错全错。"),
    ("三层兜底：retry-with-feedback → 程序救援 → 老算法兜底。",
     "第一次失败把错误反馈给模型再试；还不行就用代码强制修复。"),
], size=19)

# 关键数据 highlight
add_rect(s, Inches(0.6), Inches(5.1), Inches(12.13), Inches(1.6), COLOR_BG_SOFT)
add_text(s, Inches(0.85), Inches(5.25), Inches(12), Inches(0.5),
         "关键结果：24 个视频 100% 切分成功", size=20, bold=True, color=COLOR_ACCENT)
add_text(s, Inches(0.85), Inches(5.8), Inches(12), Inches(0.4),
         "● 22% 第一次就过    ● 67% retry 救回    ● 11% 程序救回",
         size=16, color=COLOR_TEXT)
add_text(s, Inches(0.85), Inches(6.25), Inches(12), Inches(0.4),
         "● 老算法 TextTiling fallback：0 次触发（说明前面三层够用）",
         size=16, color=COLOR_TEXT)
add_footer(s)


# ============================================================
# Page 6 - 方法 ② ASR 后处理
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 6, TOTAL, "四、核心方法 ② · ASR 后处理 与 健壮性防御")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "语音识别经常出错，靠后处理把脏数据洗干净",
         size=18, bold=True, color=COLOR_PRIMARY)

add_bullets(s, Inches(0.6), Inches(2.0), Inches(12), Inches(4.5), [
    ("术语字典自动注入。",
     "从视频标题/简介里抓专业词喂给语音模型，避免认错（如「哲学家」不会被听成「哲学者」）。"),
    ("LCP 重复段去重。",
     "把字幕里「复读机」段落合并掉——语音模型偶尔抽风，把同一句重复几十遍。"),
    ("三层健壮性防御：解码约束 + 幻觉门 + 增量落盘。",
     "长视频跑到 99% 时崩溃不会丢全部进度，类似数据库的 WAL 思想。"),
    ("ASR 置信度可视化。",
     "模型自己也不确信的字会被标 [?]，让用户能识别可疑内容。"),
], size=19)
add_footer(s)


# ============================================================
# Page 7 - 方法 ③ 学习场景结构 + 分类模板分发
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 7, TOTAL, "四、核心方法 ③ · 学习场景结构 与 模板分发")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "为不同类型视频套不同模板，学习类有 5 件套",
         size=18, bold=True, color=COLOR_PRIMARY)

# 上半：5 件套（卡片）
add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.4),
         "学习模板 5 件套", size=18, bold=True, color=COLOR_TEXT)
cards = [("目录\nTOC", "可点击跳转"),
         ("摘要卡", "全文一句话"),
         ("知识点\n速览", "每章核心点"),
         ("术语表", "首次出现链接"),
         ("章末\n小结", "抽取式总结")]
card_w, card_h = Inches(2.35), Inches(1.6)
card_gap = Inches(0.15)
total_w = card_w * 5 + card_gap * 4
start_x = (SW - total_w) / 2
for i, (head, desc) in enumerate(cards):
    x = start_x + (card_w + card_gap) * i
    y = Inches(2.4)
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    shp.fill.solid(); shp.fill.fore_color.rgb = COLOR_BG_SOFT
    shp.line.color.rgb = COLOR_PRIMARY
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = head
    r.font.name = FONT_CN; r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = COLOR_PRIMARY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc
    r2.font.name = FONT_CN; r2.font.size = Pt(11); r2.font.color.rgb = COLOR_MUTED

# 下半：分类与分发
add_bullets(s, Inches(0.6), Inches(4.4), Inches(12), Inches(2.5), [
    ("视频分类：4 类启发式分类器（教学 / 科普 / Vlog / 演讲）。",
     "用 ASR 文本特征 + metadata 自动判断，24 个视频准确率 24/24。"),
    ("模板分发：不同类视频套不同模板。",
     "Vlog 视频不需要「知识点速览」，自动换成探店风格摘要。"),
], size=18)
add_footer(s)


# ============================================================
# Page 8 - 方法 ④ 多模态视觉融合
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 8, TOTAL, "四、核心方法 ④ · 多模态视觉信号融合")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "不光看字幕，也看画面——但视觉只能当辅助",
         size=18, bold=True, color=COLOR_PRIMARY)

add_bullets(s, Inches(0.6), Inches(2.0), Inches(12), Inches(3.5), [
    ("把视频画面用 CLIP 编码，作为章节切分的辅助信号。",
     "画面跳变可能意味着章节切换。"),
    ("关键发现：视觉信号权重 α = 0.3 最佳。",
     "PPT 教学里 slide 翻页太频繁，不能让画面主导，否则会过切。"),
    ("长视频额外启用 VLM 看图模型。",
     "Qwen2.5-VL 直接描述画面内容，作为大模型的额外参考。"),
], size=19)

# 数据 highlight
add_rect(s, Inches(0.6), Inches(5.5), Inches(12.13), Inches(1.2), COLOR_BG_SOFT)
add_text(s, Inches(0.85), Inches(5.7), Inches(12), Inches(0.5),
         "多模态消融实验：9 视频中 3 个 attempts 数下降",
         size=18, bold=True, color=COLOR_ACCENT)
add_text(s, Inches(0.85), Inches(6.2), Inches(12), Inches(0.4),
         "说明视觉信号确实让大模型「少瞎想」，但 PPT 域要小心过切风险。",
         size=14, color=COLOR_MUTED)
add_footer(s)


# ============================================================
# Page 9 - 实验：Benchmark 与主结果
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 9, TOTAL, "五、实验 · Benchmark 与主结果")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "32 视频跨域 benchmark，关键指标用数字说话",
         size=18, bold=True, color=COLOR_PRIMARY)

# 数据卡片 3 个
metrics = [
    ("王道 OS p37", "0.50 → 1.00", "F1@1 翻倍", "去重救援"),
    ("计网 p38", "0.25 → 0.75", "严格 F1 +0.5", "去重 + 切分"),
    ("LLM 切分", "100%", "24 视频覆盖率", "三层兜底"),
]
card_w = Inches(3.95)
card_h = Inches(2.4)
gap = Inches(0.2)
total_w = card_w * 3 + gap * 2
start_x = (SW - total_w) / 2

for i, (title, num, sub, tag) in enumerate(metrics):
    x = start_x + (card_w + gap) * i
    y = Inches(2.0)
    add_rect(s, x, y, card_w, card_h, COLOR_BG_SOFT)
    add_text(s, x + Inches(0.2), y + Inches(0.15), card_w, Inches(0.4),
             title, size=15, bold=True, color=COLOR_MUTED)
    add_text(s, x + Inches(0.2), y + Inches(0.65), card_w, Inches(1.0),
             num, size=40, bold=True, color=COLOR_ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(1.6), card_w, Inches(0.4),
             sub, size=14, color=COLOR_TEXT)
    add_text(s, x + Inches(0.2), y + Inches(1.95), card_w, Inches(0.4),
             "→ " + tag, size=13, color=COLOR_PRIMARY)

# 补充说明
add_bullets(s, Inches(0.6), Inches(4.7), Inches(12), Inches(2.0), [
    ("Benchmark 规模：32 个视频手工标注（学习类 26 + 实拍类 6）。",
     ""),
    ("章节切分 LLM 路径覆盖：22% 一次过 + 67% retry + 11% 程序救回 = 100%。",
     ""),
    ("F1@1 是 ±1 chunk 容差指标，比严格 F1 更适合短视频评估。",
     ""),
], size=16)
add_footer(s)


# ============================================================
# Page 10 - 实验：消融与跨域泛化
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 10, TOTAL, "五、实验 · 消融与跨域泛化")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "从 10 视频扩到 32 视频，架构在不同域都稳定",
         size=18, bold=True, color=COLOR_PRIMARY)

add_bullets(s, Inches(0.6), Inches(2.0), Inches(12), Inches(4.5), [
    ("α sweep 消融：10 视频 × 2 chunker，α = 0.3 是稳健甜点。",
     "视觉权重过高会让 PPT 视频过切，过低则失去视觉信息。"),
    ("24 视频架构泛化：扩样本不退化。",
     "证明系统不是只在小数据集上调通的过拟合方案。"),
    ("中英双语支持：英文教学视频 (EH5jx5qPabU) 跑通。",
     "句号本土化、Qwen 出英文、wrap-up 大小写三处适配。"),
    ("多模态 ablation：mm 路径让 9 视频中 3 个的 LLM attempts 数减少。",
     "视觉信号确实在帮大模型做对的判断。"),
    ("跨域：教学、科普、Vlog、英文都验过。",
     "26 学习类 + 6 实拍类全部跑通。"),
], size=18)
add_footer(s)


# ============================================================
# Page 11 - 案例 ① 王道 OS p37（截图位）
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 11, TOTAL, "六、案例 ① · 王道操作系统 p37")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "ASR 卡片回路污染章节切分，去重后修复",
         size=18, bold=True, color=COLOR_PRIMARY)

# 左：问题描述
add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
         "问题 与 修复", size=18, bold=True, color=COLOR_TEXT)
add_bullets(s, Inches(0.6), Inches(2.5), Inches(6), Inches(3.5), [
    ("问题：「哲学家进餐」段被语音模型复读了 8 遍。", ""),
    ("结果：原始章节切分被关键词频次干扰彻底乱掉。", ""),
    ("方法：LCP 共同前缀检测连续重复段并合并时间戳。", ""),
    ("效果：F1@1 从 0.50 提升到 1.00。", ""),
], size=15)

# 右：截图占位
add_fitted_picture(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.7),
                   "screenshots/01_os_chapters.png")
add_footer(s)


# ============================================================
# Page 12 - 案例 ② vlog hallucinate（截图位）
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 12, TOTAL, "六、案例 ② · Vlog Abstract 字面 Hallucinate 修复")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "大模型「自由发挥」写出与原视频无关的内容",
         size=18, bold=True, color=COLOR_PRIMARY)

add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
         "问题 与 修复", size=18, bold=True, color=COLOR_TEXT)
add_bullets(s, Inches(0.6), Inches(2.5), Inches(6), Inches(3.5), [
    ("问题：BV1q6 日料探店，13/16 章 abstract 瞎编。", ""),
    ("例：「心理线」被写成「股市心理」、「皇上」被写成「清朝」。", ""),
    ("原因：模型只看大纲不看原话，开始自由发挥。", ""),
    ("修复：prompt 里塞一小段原字幕作参考。", ""),
    ("效果：3 个 Vlog 视频 21/21 章 0 hallucinate。", ""),
], size=15)

add_fitted_picture(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.7),
                   "screenshots/02_vlog_abstract.png")
add_footer(s)


# ============================================================
# Page 13 - 案例 ③ Web 前端（截图位）
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 13, TOTAL, "六、案例 ③ · Web 前端展示")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "Next.js 16 + Plyr 播放器 + 双语切换",
         size=18, bold=True, color=COLOR_PRIMARY)

add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
         "前端功能", size=18, bold=True, color=COLOR_TEXT)
add_bullets(s, Inches(0.6), Inches(2.5), Inches(6), Inches(3.5), [
    ("Plyr 视频播放器嵌入 + 章节按钮直接跳转。", ""),
    ("章节卡片显示当前播放进度条与关键词标签。", ""),
    ("Apple 风格深色 UI + 流动粒子背景。", ""),
    ("章节摘要、知识点、术语表三栏可滚动浏览。", ""),
    ("整套前端 SSG 静态化，加载快、可离线。", ""),
], size=15)

add_fitted_picture(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.7),
                   "screenshots/03_web_overall.png")
add_footer(s)


# ============================================================
# Page 14 - 结论 & 未来工作
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 14, TOTAL, "七、结论 与 未来工作")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "完成端到端系统 + 跨域 benchmark + 算法/工程双重创新",
         size=18, bold=True, color=COLOR_PRIMARY)

# 左：完成
add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
         "已完成", size=18, bold=True, color=COLOR_ACCENT)
add_bullets(s, Inches(0.6), Inches(2.5), Inches(6), Inches(3.5), [
    "32 视频跨域 benchmark 手工标注",
    "LLM 章节切分 100% 覆盖率",
    "学习场景 5 件套 md 模板",
    "4 类启发式分类 24/24 准",
    "中英双语支持",
    "Next.js 16 web 前端 + 5 demo",
], size=15)

# 右：未来 + 局限
add_text(s, Inches(7.0), Inches(2.0), Inches(6), Inches(0.4),
         "局限 与 未来工作", size=18, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(7.0), Inches(2.5), Inches(6), Inches(3.5), [
    "小模型在 30+ chunks 极长视频上偶发失效",
    "PPT slide cue 跨视频校准困难",
    "扩 talk 类视频（讲座、播客）",
    "研究更小更快的章节切分模型",
    "用户主观打分的大规模评估",
], size=15)

# 底部致谢
add_rect(s, 0, Inches(6.7), SW, Inches(0.8), COLOR_PRIMARY)
add_text(s, Inches(0.6), Inches(6.88), Inches(12), Inches(0.5),
         "感谢聆听 · Q & A",
         size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# 保存
# ============================================================
out = "notegen_presentation.pptx"
prs.save(out)
print(f"OK: {out}  ({len(prs.slides)} slides)")
